import collections 
import collections.abc
from datetime import datetime
from pptx.enum.text import PP_ALIGN
from pptx import Presentation
from pptx.util import Inches
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.util import Pt
datetime_now = datetime.now()
currencies_ppt = Presentation('oldppt.pptx')

# Slide 1 
slide1 = currencies_ppt.slides[0]
for shape in slide1.shapes:
    # print(shape.shape_id)
    # print(shape.shape_type)
    if shape.shape_id == 7 :
        text_frame = shape.text_frame
        text_frame.clear()
        p = text_frame.paragraphs[0]
        run = p.add_run()
        run.text = "Automating x\nWeekly Report\nWith Python\nStanley"
        font = run.font
        font.name = 'Segoe UI'
        font.size = Pt(32)
        font.bold = True
        font.italic = None  # cause value to be inherited from theme
        font.color.rgb = RGBColor(0x10, 0x0c, 0x08)

# Slide 2
slide2 = currencies_ppt.slides[1]
for shape in slide2.shapes:
    # print(shape.shape_id)
    # print(shape.shape_type)
    if shape.shape_id == 7 :
        text_frame = shape.text_frame
        p = text_frame.paragraphs[0]
        text_frame.clear()  
        run = p.add_run()
        run.text = "Agenda"
        font = run.font
        font.name = 'Century Gothic'
        font.size = Pt(22)
        font.bold = True
        font.color.rgb = RGBColor(0x10, 0x0c, 0x08)
    if shape.shape_id == 2 :
        text_frame = shape.text_frame

        p1 = text_frame.paragraphs[0]
        p1.text = "Dataset summary"
        font = p1.font
        font.name = 'Century Gothic'
        font.size = Pt(18)
        font.color.rgb = RGBColor(0x10, 0x0c, 0x08)

        p1 = text_frame.paragraphs[1]
        p1.text = "Model v2 development"
        font = p1.font
        font.name = 'Century Gothic'
        font.size = Pt(18)
        font.color.rgb = RGBColor(0x10, 0x0c, 0x08)

        p1 = text_frame.paragraphs[2]
        p1.text = "Methodology"
        font = p1.font
        font.name = 'Century Gothic'
        font.size = Pt(18)
        font.color.rgb = RGBColor(0x10, 0x0c, 0x08)

        p1 = text_frame.paragraphs[3]
        p1.text = "Model selection"
        font = p1.font
        font.name = 'Century Gothic'
        font.size = Pt(18)
        font.color.rgb = RGBColor(0x10, 0x0c, 0x08)

        p1 = text_frame.paragraphs[4]
        p1.text = "Feature importance - SHAP"
        font = p1.font
        font.name = 'Century Gothic'
        font.size = Pt(18)
        font.color.rgb = RGBColor(0x10, 0x0c, 0x08)

        p1 = text_frame.paragraphs[5]
        p1.text = "Performance comparison: Model v1 vs Model v2"
        font = p1.font
        font.name = 'Century Gothic'
        font.size = Pt(18)
        font.color.rgb = RGBColor(0x10, 0x0c, 0x08)

# Slide 3
slide = currencies_ppt.slides[2]
shapes = slide.shapes
for shape in slide.shapes:
    # print(shape.shape_id)
    # print(shape.shape_type)
    if shape.shape_id == 7 :
        text_frame = shape.text_frame
        p1 = text_frame.paragraphs[0]
        p1.text = "Dataset Summary"
        font = p1.font
        font.name = 'Century Gothic'
        font.size = Pt(22)
        font.bold = True
        font.color.rgb = RGBColor(0x10, 0x0c, 0x08)
    if shape.shape_id == 9 :
        text_frame = shape.text_frame
        p1 = text_frame.paragraphs[0]
        p1.text = "Development sample: (202101 to 202104) and 202107"
        font = p1.font
        font.name = 'Century Gothic'
        font.size = Pt(14)
        font.color.rgb = RGBColor(0x10, 0x0c, 0x08)

        p1 = text_frame.paragraphs[1]
        p1.text = "Covid testing sample: 202105 to 202106"
        font = p1.font
        font.name = 'Century Gothic'
        font.size = Pt(14)
        font.color.rgb = RGBColor(0x10, 0x0c, 0x08)

        p1 = text_frame.paragraphs[2]
        p1.text = "OOT sample: 202108 to 202109"
        font = p1.font
        font.name = 'Century Gothic'
        font.size = Pt(14)
        font.color.rgb = RGBColor(0x10, 0x0c, 0x08)
    if shape.shape_id == 6 :
        shapes.element.remove(shape.element)
        img_path = 'slide3_1.jpg'
        top_pic = slide.shapes.add_picture(img_path, Inches(0.40), Inches(2.50), width=Inches(5.30))

# Slide 5
slide = currencies_ppt.slides[3]
shapes = slide.shapes
for shape in slide.shapes:
    # print(shape.shape_id)
    # print(shape.shape_type)
    if shape.shape_id == 7 :
        text_frame = shape.text_frame
        p1 = text_frame.paragraphs[0]
        p1.text = "Model v2 Development – methodology"
        font = p1.font
        font.name = 'Century Gothic'
        font.size = Pt(22)
        font.bold = True
        font.color.rgb = RGBColor(0x10, 0x0c, 0x08)
    if shape.shape_id == 5 :
        text_frame = shape.text_frame
        p1 = text_frame.paragraphs[0]
        p1.text = "Time based cross-validation"
        font = p1.font
        font.name = 'Century Gothic'
        font.size = Pt(14)
        font.color.rgb = RGBColor(0x10, 0x0c, 0x08)

        p1 = text_frame.paragraphs[8]
        p1.text = "Each fold:"
        font = p1.font
        font.name = 'Century Gothic'
        font.size = Pt(14)
        font.color.rgb = RGBColor(0x10, 0x0c, 0x08)

        p1 = text_frame.paragraphs[9]
        p1.text = "Shifts forward 1 weekstart"
        font = p1.font
        font.name = 'Century Gothic'
        font.size = Pt(14)
        font.color.rgb = RGBColor(0x10, 0x0c, 0x08)

        p1 = text_frame.paragraphs[10]
        p1.text = "Has n week_train (blue)"
        font = p1.font
        font.name = 'Century Gothic'
        font.size = Pt(14)
        font.color.rgb = RGBColor(0x10, 0x0c, 0x08)

        p1 = text_frame.paragraphs[11]
        p1.text = "Has m week_test (orange)"
        font = p1.font
        font.name = 'Century Gothic'
        font.size = Pt(14)
        font.color.rgb = RGBColor(0x10, 0x0c, 0x08)

        p1 = text_frame.paragraphs[12]
        p1.text = "Last n weekstart will be used for final model"
        font = p1.font
        font.name = 'Century Gothic'
        font.size = Pt(14)
        font.color.rgb = RGBColor(0x10, 0x0c, 0x08)

        p1 = text_frame.paragraphs[13]
        p1.text = "Only using development sample type"
        font = p1.font
        font.name = 'Century Gothic'
        font.size = Pt(14)
        font.color.rgb = RGBColor(0x10, 0x0c, 0x08)

        p1 = text_frame.paragraphs[14]
        p1.text = "Weekstart: [2020-12-28 to 2021-04-26] and [2021-06-28 to 2021-07-26]"
        font = p1.font
        font.name = 'Century Gothic'
        font.size = Pt(14)
        font.color.rgb = RGBColor(0x10, 0x0c, 0x08)
    if shape.shape_id == 6 :
        shapes.element.remove(shape.element)
        img_path = 'image-20220422-071652.png'
        top_pic = slide.shapes.add_picture(img_path, Inches(0.80), Inches(1.50), width=Inches(4.2))

def iter_cells(table):
    for row in table.rows:
        for cell in row.cells:
            yield cell

# Slide 6
slide = currencies_ppt.slides[4]
shapes = slide.shapes
for shape in slide.shapes:
    # print(shape.shape_id)
    # print(shape.shape_type)
    if shape.shape_id == 7 :
        text_frame = shape.text_frame
        p1 = text_frame.paragraphs[0]
        p1.text = "Model v2 Development – model selection"
        font = p1.font
        font.name = 'Century Gothic'
        font.size = Pt(22)
        font.bold = True
        font.color.rgb = RGBColor(0x10, 0x0c, 0x08)
    if shape.shape_id == 8 :
        text_frame = shape.text_frame
        p1 = text_frame.paragraphs[0]
        p1.text = "Considering the trade-off between performance and number of features, we recommend model_8 as the selected model for production"
        font = p1.font
        font.name = 'Century Gothic'
        font.size = Pt(18)
        font.color.rgb = RGBColor(0x10, 0x0c, 0x08)
    if shape.shape_id == 2 :
        table = shape.table
        cell = table.cell(3, 4)
        cell.text = '0.493'
        cell = table.cell(0, 1)
        cell.text = 'Number of features'
        for cell in iter_cells(table):
            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.CENTER
                for run in paragraph.runs:
                    run.font.size = Pt(14)
                    run.font.name = 'Century Gothic'
                    run.font.color.rgb = RGBColor(0x10, 0x0c, 0x08)

# Slide 11
slide = currencies_ppt.slides[5]
shapes = slide.shapes
for shape in slide.shapes:
    # print(shape.shape_id)
    # print(shape.shape_type)
    if shape.shape_id == 7 :
        text_frame = shape.text_frame
        p1 = text_frame.paragraphs[0]
        p1.text = "Performance comparison – GINI per month"
        font = p1.font
        font.name = 'Century Gothic'
        font.size = Pt(22)
        font.bold = True
        font.color.rgb = RGBColor(0x10, 0x0c, 0x08)
    if shape.shape_id == 8 :
        text_frame = shape.text_frame
        p1 = text_frame.paragraphs[0]
        p1.text = "No calibration vs model_1 (baseline) vs model_8"
        font = p1.font
        font.name = 'Century Gothic'
        font.size = Pt(14)
        font.bold = True
        font.underline = True
        font.color.rgb = RGBColor(0x10, 0x0c, 0x08)
    if shape.shape_id == 5 :
        shapes.element.remove(shape.element)
        img_path = 'slide11_1.jpg'
        top_pic = slide.shapes.add_picture(img_path, Inches(0.35), Inches(1.30), width=Inches(6.45))
    if shape.shape_id == 2 :
        shapes.element.remove(shape.element)
        img_path = 'slide11_2.jpg'
        top_pic = slide.shapes.add_picture(img_path, Inches(0.35), Inches(3.20), width=Inches(7.74))

currencies_ppt.save('newppt.pptx')